"""Format matrix: generate a real file in every supported format, parse it,
assert the contract text survived the round trip.

Not a plumbing test. Each case defends the claim "you can drop your documents
in": if a format silently yields empty or truncated text, source_span
verification fails on values that genuinely are in the document.
"""
import json
import os
import shutil
import subprocess
import sys
import tempfile
import zipfile

import ingest

# Long enough to clear the MIN_CHARS text-layer floor.
BODY = (
    "MASTER SERVICES AGREEMENT between Meridian Holdings LLC and Northgate "
    "Legal Services. This Agreement is effective January 1, 2026 and shall "
    "remain in force until March 31, 2027. Section 14.3: this Agreement shall "
    "automatically renew for successive twelve-month terms unless either party "
    "gives written notice at least sixty (60) days prior to expiry. Fees are "
    "USD 120,000 per annum, payable quarterly. Governing law: Delaware."
)
NEEDLE = "sixty (60) days"
PASS, FAIL = [], []


def check(label, path, needle=NEEDLE, expect_fmt=None, expect_conv=None):
    try:
        doc = ingest.parse(path)
    except ingest.Unsupported as exc:
        FAIL.append(f"{label}: REFUSED - {exc}")
        return None
    except Exception as exc:
        FAIL.append(f"{label}: {type(exc).__name__} - {exc}")
        return None
    if needle and needle not in doc.text:
        FAIL.append(f"{label}: needle missing from {len(doc.text)} chars")
        return None
    if expect_fmt and doc.fmt != expect_fmt:
        FAIL.append(f"{label}: fmt={doc.fmt} expected {expect_fmt}")
        return None
    if expect_conv and doc.converted_via != expect_conv:
        FAIL.append(f"{label}: converted_via={doc.converted_via}")
        return None
    PASS.append(f"{label:<26} fmt={doc.fmt:<6} chars={len(doc.text):<6}"
                f"pages={len(doc.pages)}"
                + (f" via={doc.converted_via}" if doc.converted_via else ""))
    return doc


tmp = tempfile.mkdtemp(prefix="fmtmatrix-")


def p(name):
    return os.path.join(tmp, name)


# ---- plain text family
for name, fmt in (("c.txt", "txt"), ("c.log", "txt")):
    open(p(name), "w").write(BODY)
    check(name, p(name), expect_fmt=fmt)

open(p("c.md"), "w").write(f"# Agreement\n\n## Section 14.3\n\n{BODY}\n")
check("c.md", p("c.md"), expect_fmt="md")

# ---- pdf (with a real text layer, two pages)
import fitz
pdf = fitz.open()
half = len(BODY) // 2
for chunk in (BODY[:half], BODY[half:]):
    page = pdf.new_page()
    page.insert_textbox(fitz.Rect(50, 50, 550, 750), chunk, fontsize=11)
pdf.save(p("c.pdf"))
pdf.close()
d = check("c.pdf", p("c.pdf"), needle=None, expect_fmt="pdf")
if d and len(d.pages) != 2:
    FAIL.append(f"c.pdf: expected 2 pages, got {len(d.pages)}")

# ---- docx: body, TABLE, and header. The table is the point: fee schedules
# ---- live in tables and a paragraph-only reader drops them.
from docx import Document
doc = Document()
doc.add_paragraph(BODY)
table = doc.add_table(rows=2, cols=2)
table.cell(0, 0).text = "Annual fee"
table.cell(0, 1).text = "USD 120,000 escalating at CPI plus 3%"
table.cell(1, 0).text = "Renewal term"
table.cell(1, 1).text = "twelve (12) months"
doc.sections[0].header.paragraphs[0].text = "CONFIDENTIAL - Meridian MSA"
doc.save(p("c.docx"))
d = check("c.docx", p("c.docx"), expect_fmt="docx")
if d:
    for must in ("CPI plus 3%", "CONFIDENTIAL - Meridian MSA", "twelve (12) months"):
        if must not in d.text:
            FAIL.append(f"c.docx: lost {must!r} (table/header extraction broken)")

# ---- rtf
open(p("c.rtf"), "w").write(
    r"{\rtf1\ansi\deff0 " + BODY.replace("\n", r"\par ") + "}")
check("c.rtf", p("c.rtf"), expect_fmt="rtf")

# ---- odt / ods / odp via odfpy
from odf.opendocument import OpenDocumentText
from odf.text import P
odt = OpenDocumentText()
odt.text.addElement(P(text=BODY))
odt.save(p("c.odt"))
check("c.odt", p("c.odt"), expect_fmt="odt")

# ---- xlsx
import openpyxl
wb = openpyxl.Workbook()
ws = wb.active
ws.title = "Terms"
ws["A1"], ws["B1"] = "Clause", "Text"
ws["A2"], ws["B2"] = "14.3", BODY
wb.create_sheet("Fees")["A1"] = "USD 120,000 per annum"
wb.save(p("c.xlsx"))
d = check("c.xlsx", p("c.xlsx"), expect_fmt="xlsx")
if d and "USD 120,000 per annum" not in d.text:
    FAIL.append("c.xlsx: second sheet not read")

# ---- pptx
from pptx import Presentation
from pptx.util import Inches
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9),
                         Inches(5)).text_frame.text = BODY
prs.save(p("c.pptx"))
check("c.pptx", p("c.pptx"), expect_fmt="pptx")

# ---- html
open(p("c.html"), "w").write(
    f"<html><head><style>b{{color:red}}</style>"
    f"<script>var x=1</script></head><body><h1>MSA</h1><p>{BODY}</p>"
    f"</body></html>")
d = check("c.html", p("c.html"), expect_fmt="html")
if d and ("var x=1" in d.text or "color:red" in d.text):
    FAIL.append("c.html: script/style leaked into text")

# ---- csv / tsv
open(p("c.csv"), "w").write(f'clause,text\n"14.3","{BODY}"\n')
check("c.csv", p("c.csv"), expect_fmt="csv")
open(p("c.tsv"), "w").write(f"clause\ttext\n14.3\t{BODY}\n")
check("c.tsv", p("c.tsv"), expect_fmt="tsv")

# ---- json
json.dump({"contract": {"clauses": [{"id": "14.3", "body": BODY}]}},
          open(p("c.json"), "w"))
check("c.json", p("c.json"), expect_fmt="json")

# ---- eml with a docx attachment: attachment must be parsed recursively
from email.message import EmailMessage
msg = EmailMessage()
msg["From"], msg["To"] = "counsel@firm.example", "partner@firm.example"
msg["Subject"] = "Meridian MSA for review"
msg.set_content("Please review the attached agreement before the renewal date.")
with open(p("c.docx"), "rb") as fh:
    msg.add_attachment(fh.read(), maintype="application",
                       subtype="vnd.openxmlformats-officedocument"
                               ".wordprocessingml.document",
                       filename="meridian.docx")
open(p("c.eml"), "wb").write(msg.as_bytes())
d = check("c.eml", p("c.eml"), expect_fmt="eml")
if d:
    if "meridian.docx" not in " ".join(d.notes):
        FAIL.append("c.eml: attachment not parsed")
    if "partner@firm.example" not in d.text:
        FAIL.append("c.eml: headers missing")

# ---- zip containing several members
with zipfile.ZipFile(p("c.zip"), "w") as zf:
    zf.write(p("c.docx"), "contracts/meridian.docx")
    zf.write(p("c.txt"), "contracts/side-letter.txt")
    zf.writestr("readme.bin", b"\x00\x01binary junk")
d = check("c.zip", p("c.zip"), expect_fmt="zip")
if d:
    notes = " ".join(d.notes)
    if "meridian.docx" not in notes or "side-letter.txt" not in notes:
        FAIL.append("c.zip: members not parsed")
    if "readme.bin" not in notes:
        FAIL.append("c.zip: unsupported member not reported")

# ---- extensionless file, identified by content sniffing
shutil.copy(p("c.pdf"), p("mystery"))
check("mystery (sniffed pdf)", p("mystery"), needle=None, expect_fmt="pdf")

# ---- Tier 2: legacy .doc via LibreOffice
if shutil.which("soffice"):
    subprocess.run(["soffice", "--headless", "--norestore", "--convert-to",
                    "doc", "--outdir", tmp, p("c.txt")],
                   check=True, timeout=180,
                   stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    legacy = p("c.doc")
    if os.path.exists(legacy):
        check("c.doc (LibreOffice)", legacy, expect_fmt="doc",
              expect_conv="libreoffice")
    else:
        FAIL.append("c.doc: LibreOffice produced no .doc to test")
else:
    FAIL.append("soffice missing - Tier 2 untestable")

# ---- Tier 3: refusals must be loud, with a reason
def must_refuse(label, path, expect_substr):
    try:
        ingest.parse(path)
    except ingest.Unsupported as exc:
        if expect_substr not in str(exc):
            FAIL.append(f"{label}: refused but reason lacked {expect_substr!r}: {exc}")
        else:
            PASS.append(f"{label:<26} REFUSED: {exc}")
        return
    FAIL.append(f"{label}: ACCEPTED - should have been refused")


# image-only pdf: real PDF, zero text layer. The dangerous case.
blank = fitz.open()
blank.new_page()
blank.save(p("scan.pdf"))
blank.close()
must_refuse("scan.pdf (no text layer)", p("scan.pdf"), "OCR is out of scope")

open(p("tiny.txt"), "w").write("too short")
must_refuse("tiny.txt (under floor)", p("tiny.txt"), "no readable text")

open(p("c.xyz"), "w").write(BODY)
must_refuse("c.xyz (unknown type)", p("c.xyz"), "unsupported file type")

# ---- offsets: every parsed doc must be able to locate its own text
d = ingest.parse(p("c.pdf"))
hit = d.locate("sixty (60) days")
if not hit:
    # PDF text layers wrap lines; locate() is whitespace-tolerant, so this
    # must still resolve.
    FAIL.append("c.pdf: locate() could not find a span in its own text")
else:
    s, e, pg = hit
    PASS.append(f"{'c.pdf locate()':<26} chars {s}-{e} on page {pg}")

shutil.rmtree(tmp, ignore_errors=True)

for line in PASS:
    print("  ok   " + line)
for line in FAIL:
    print("  FAIL " + line)
print(f"\n{len(PASS)} passed, {len(FAIL)} failed")
if FAIL:
    sys.exit(1)
print("ALL FORMAT TESTS PASSED")

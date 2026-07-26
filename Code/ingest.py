"""Document ingest: one registry, many formats, deterministic output.

Contract with the rest of the system:
  * parse(path) -> ParsedDoc            the only entry point
  * ParsedDoc.text is THE text every validator checks source_spans against.
  * Page.char_start lets any (start,end) offset into .text be resolved back to
    a page number -- that is what makes a quote clickable in the UI (D-B).

Tier 1  native python parsers          (pdf docx txt md rtf odt ods odp xlsx
                                        pptx html csv tsv eml json zip)
Tier 2  LibreOffice headless fallback  (doc xls ppt wpd ...) -> Tier 1
Tier 3  refused with a reason          (no text layer / image-only / empty)

Tier 3 exists because an empty extraction is the most dangerous possible
input: with no text, every source_span check is vacuous and a fabricated
extraction would sail through validation. We refuse instead.
"""
from __future__ import annotations

import csv
import io
import json
import os
import re
import subprocess
import tempfile
import zipfile
from dataclasses import dataclass, field
from email import policy as email_policy
from email.parser import BytesParser

# Minimum characters before we believe a document actually has a text layer.
MIN_CHARS = int(os.environ.get("LEDGER_MIN_CHARS", "200"))
SOFFICE_TIMEOUT = int(os.environ.get("LEDGER_SOFFICE_TIMEOUT", "120"))


class Unsupported(Exception):
    """Tier 3: we will not guess. Carries a human-readable reason."""


@dataclass
class Page:
    number: int          # 1-based
    text: str
    char_start: int      # offset of this page's text within ParsedDoc.text


@dataclass
class Passage:
    """Shared with the RAG lane. Coordinates are offsets into ParsedDoc.text."""
    text: str
    contract_id: int
    page: int
    char_start: int
    char_end: int
    score: float = 0.0


@dataclass
class ParsedDoc:
    text: str
    pages: list[Page] = field(default_factory=list)
    fmt: str = ""
    converted_via: str | None = None
    notes: list[str] = field(default_factory=list)

    def page_of(self, char_offset: int) -> int:
        """Resolve a character offset to a 1-based page number."""
        page = 1
        for p in self.pages:
            if p.char_start > char_offset:
                break
            page = p.number
        return page

    def locate(self, span: str) -> tuple[int, int, int] | None:
        """Find a verbatim span. Returns (start, end, page) or None.

        Exact match first, then a whitespace-tolerant regex, because PDF text
        layers break lines inside sentences. Never fuzzy beyond whitespace --
        V1 must stay a real check.
        """
        if not span or not span.strip():
            return None
        i = self.text.find(span)
        if i != -1:
            return i, i + len(span), self.page_of(i)
        pattern = r"\s+".join(re.escape(w) for w in span.split())
        m = re.search(pattern, self.text, re.I)
        if m:
            return m.start(), m.end(), self.page_of(m.start())
        return None


def _decode(raw: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _paged(chunks: list[str], sep: str = "\n\n") -> tuple[str, list[Page]]:
    """Assemble numbered chunks into one text plus page offsets."""
    pages: list[Page] = []
    parts: list[str] = []
    cursor = 0
    for n, chunk in enumerate(chunks, start=1):
        pages.append(Page(number=n, text=chunk, char_start=cursor))
        parts.append(chunk)
        cursor += len(chunk) + len(sep)
    return sep.join(parts), pages


# --------------------------------------------------------------- Tier 1

def _pdf(path: str) -> ParsedDoc:
    """Per-page text with offsets.

    MuPDF writes font complaints ("unknown cid font type") straight to stderr,
    one per glyph run. A real 40-page lease produced 400+ lines of it, which
    buried the actual result and made a working parse look like a failure.
    They are advisory -- text still extracts -- so the messages are collected
    into doc.notes instead of the console.
    """
    import fitz
    if hasattr(fitz, "TOOLS") and hasattr(fitz.TOOLS, "mupdf_display_errors"):
        fitz.TOOLS.mupdf_display_errors(False)
    with fitz.open(path) as doc:
        chunks = [p.get_text() for p in doc]
    text, pages = _paged(chunks)
    out = ParsedDoc(text=text, pages=pages, fmt="pdf")
    warned = 0
    if hasattr(fitz, "TOOLS") and hasattr(fitz.TOOLS, "mupdf_warnings"):
        warned = len((fitz.TOOLS.mupdf_warnings() or "").splitlines())
    if warned:
        out.notes.append(f"{warned} MuPDF font warning(s); text extracted "
                         "normally, but embedded fonts are non-standard")
    return out


def _docx(path: str) -> ParsedDoc:
    """Paragraphs AND tables AND headers/footers.

    Not thoroughness for its own sake: fee schedules live in tables and
    renewal clauses hide in late sections. A paragraph-only reader drops them,
    and then a correct source_span fails V1 because the text was never read.
    """
    from docx import Document
    from docx.document import Document as _Doc
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table, _Cell
    from docx.text.paragraph import Paragraph

    def walk(parent):
        if isinstance(parent, _Doc):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        else:
            return
        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent).text
            elif isinstance(child, CT_Tbl):
                table = Table(child, parent)
                for row in table.rows:
                    cells = []
                    for cell in row.cells:
                        cells.append(" ".join(t for t in walk(cell) if t))
                    yield " | ".join(cells)

    doc = Document(path)
    blocks = [t for t in walk(doc) if t and t.strip()]
    for section in doc.sections:
        for container in (section.header, section.footer):
            for para in container.paragraphs:
                if para.text.strip():
                    blocks.append(para.text)
    return ParsedDoc(text="\n".join(blocks), fmt="docx")


def _txt(path: str) -> ParsedDoc:
    with open(path, "rb") as fh:
        return ParsedDoc(text=_decode(fh.read()), fmt="txt")


def _md(path: str) -> ParsedDoc:
    with open(path, "rb") as fh:
        raw = _decode(fh.read())
    stripped = re.sub(r"^#{1,6}\s*|^\s*[-*+]\s+|`{1,3}|\*{1,2}|_{1,2}",
                      "", raw, flags=re.M)
    return ParsedDoc(text=stripped, fmt="md")


def _rtf(path: str) -> ParsedDoc:
    from striprtf.striprtf import rtf_to_text
    with open(path, "rb") as fh:
        return ParsedDoc(text=rtf_to_text(_decode(fh.read())), fmt="rtf")


def _odf(path: str) -> ParsedDoc:
    from odf import teletype, text as odftext
    from odf.opendocument import load
    doc = load(path)
    blocks = [teletype.extractText(p)
              for p in doc.getElementsByType(odftext.P)]
    ext = os.path.splitext(path)[1].lstrip(".").lower()
    return ParsedDoc(text="\n".join(b for b in blocks if b.strip()), fmt=ext)


def _xlsx(path: str) -> ParsedDoc:
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    chunks = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        chunks.append(f"[sheet: {ws.title}]\n" + "\n".join(rows))
    wb.close()
    text, pages = _paged(chunks)
    return ParsedDoc(text=text, pages=pages, fmt="xlsx")


def _pptx(path: str) -> ParsedDoc:
    from pptx import Presentation
    prs = Presentation(path)
    chunks = []
    for slide in prs.slides:
        bits = [sh.text for sh in slide.shapes
                if getattr(sh, "has_text_frame", False) and sh.text.strip()]
        chunks.append("\n".join(bits))
    text, pages = _paged(chunks)
    return ParsedDoc(text=text, pages=pages, fmt="pptx")


def _html(path: str) -> ParsedDoc:
    from bs4 import BeautifulSoup
    with open(path, "rb") as fh:
        soup = BeautifulSoup(_decode(fh.read()), "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    return ParsedDoc(text=soup.get_text("\n", strip=True), fmt="html")


def _csv(path: str) -> ParsedDoc:
    with open(path, "rb") as fh:
        raw = _decode(fh.read())
    try:
        dialect = csv.Sniffer().sniff(raw[:4096])
    except csv.Error:
        dialect = csv.excel
    rows = [" | ".join(c for c in row if c)
            for row in csv.reader(io.StringIO(raw), dialect)]
    ext = os.path.splitext(path)[1].lstrip(".").lower()
    return ParsedDoc(text="\n".join(r for r in rows if r.strip()), fmt=ext)


def _json(path: str) -> ParsedDoc:
    with open(path, "rb") as fh:
        data = json.loads(_decode(fh.read()))
    out: list[str] = []

    def harvest(node, prefix=""):
        if isinstance(node, dict):
            for k, v in node.items():
                harvest(v, f"{prefix}{k}: " if isinstance(v, str) else prefix)
        elif isinstance(node, list):
            for v in node:
                harvest(v, prefix)
        elif isinstance(node, str):
            out.append(prefix + node)
        elif node is not None:
            out.append(f"{prefix}{node}")

    harvest(data)
    return ParsedDoc(text="\n".join(out), fmt="json")


def _eml(path: str) -> ParsedDoc:
    """Body plus attachments, each attachment re-entering the registry."""
    with open(path, "rb") as fh:
        msg = BytesParser(policy=email_policy.default).parse(fh)
    head = "\n".join(f"{k}: {msg[k]}" for k in
                     ("From", "To", "Cc", "Date", "Subject") if msg[k])
    chunks = [head]
    notes: list[str] = []
    body = msg.get_body(preferencelist=("plain", "html"))
    if body is not None:
        payload = body.get_content()
        if body.get_content_type() == "text/html":
            from bs4 import BeautifulSoup
            payload = BeautifulSoup(payload, "lxml").get_text("\n", strip=True)
        chunks.append(payload)
    for part in msg.iter_attachments():
        name = part.get_filename() or "attachment"
        ext = os.path.splitext(name)[1].lower()
        if ext not in HANDLERS and ext not in SOFFICE_TARGETS:
            notes.append(f"attachment skipped (unsupported): {name}")
            continue
        with tempfile.TemporaryDirectory() as tmp:
            dest = os.path.join(tmp, os.path.basename(name))
            with open(dest, "wb") as out:
                out.write(part.get_payload(decode=True) or b"")
            try:
                sub = parse(dest, _depth=1)
            except Unsupported as exc:
                notes.append(f"attachment skipped ({exc}): {name}")
                continue
        chunks.append(f"[attachment: {name}]\n{sub.text}")
        notes.append(f"attachment parsed: {name} ({sub.fmt})")
    text, pages = _paged(chunks)
    doc = ParsedDoc(text=text, pages=pages, fmt="eml")
    doc.notes = notes
    return doc


def _zip(path: str) -> ParsedDoc:
    """Every member re-enters the registry. Members that cannot be parsed are
    recorded as notes, never silently dropped."""
    chunks, notes = [], []
    with zipfile.ZipFile(path) as zf, tempfile.TemporaryDirectory() as tmp:
        for info in zf.infolist():
            if info.is_dir():
                continue
            ext = os.path.splitext(info.filename)[1].lower()
            if ext not in HANDLERS and ext not in SOFFICE_TARGETS:
                notes.append(f"member skipped (unsupported): {info.filename}")
                continue
            dest = os.path.join(tmp, os.path.basename(info.filename))
            with zf.open(info) as src, open(dest, "wb") as out:
                out.write(src.read())
            try:
                sub = parse(dest, _depth=1)
            except Unsupported as exc:
                notes.append(f"member skipped ({exc}): {info.filename}")
                continue
            chunks.append(f"[member: {info.filename}]\n{sub.text}")
            notes.append(f"member parsed: {info.filename} ({sub.fmt})")
    if not chunks:
        raise Unsupported("zip contained no parseable documents")
    text, pages = _paged(chunks)
    doc = ParsedDoc(text=text, pages=pages, fmt="zip")
    doc.notes = notes
    return doc


HANDLERS: dict[str, callable] = {
    ".pdf": _pdf, ".docx": _docx, ".docm": _docx,
    ".txt": _txt, ".text": _txt, ".log": _txt,
    ".md": _md, ".markdown": _md,
    ".rtf": _rtf,
    ".odt": _odf, ".ods": _odf, ".odp": _odf,
    ".xlsx": _xlsx, ".xlsm": _xlsx,
    ".pptx": _pptx,
    ".html": _html, ".htm": _html, ".xhtml": _html,
    ".csv": _csv, ".tsv": _csv,
    ".json": _json,
    ".eml": _eml, ".mht": _html,
    ".zip": _zip,
}

# Tier 2: legacy formats LibreOffice converts, mapped to its --convert-to target.
SOFFICE_TARGETS: dict[str, str] = {
    ".doc": "docx", ".dot": "docx", ".wpd": "docx", ".pages": "docx",
    ".xls": "xlsx", ".xlsb": "xlsx",
    ".ppt": "pptx", ".pps": "pptx",
}

SNIFF: list[tuple[bytes, str]] = [
    (b"%PDF-", ".pdf"),
    (b"{\\rtf", ".rtf"),
    (b"\xd0\xcf\x11\xe0", ".doc"),   # OLE2 -> let LibreOffice decide
]


def _sniff(path: str) -> str | None:
    try:
        with open(path, "rb") as fh:
            head = fh.read(8)
    except OSError:
        return None
    for magic, ext in SNIFF:
        if head.startswith(magic):
            return ext
    if head.startswith(b"PK\x03\x04"):
        try:
            with zipfile.ZipFile(path) as zf:
                names = set(zf.namelist())
        except zipfile.BadZipFile:
            return None
        if "word/document.xml" in names:
            return ".docx"
        if any(n.startswith("xl/") for n in names):
            return ".xlsx"
        if any(n.startswith("ppt/") for n in names):
            return ".pptx"
        if "mimetype" in names:
            return ".odt"
        return ".zip"
    return None


def _soffice(path: str, target: str) -> str:
    """Convert via LibreOffice headless. Offline, timeout-bounded."""
    outdir = tempfile.mkdtemp(prefix="ledger-conv-")
    cmd = ["soffice", "--headless", "--norestore",
           "--convert-to", target, "--outdir", outdir, path]
    try:
        subprocess.run(cmd, check=True, timeout=SOFFICE_TIMEOUT,
                       stdout=subprocess.DEVNULL, stderr=subprocess.PIPE)
    except FileNotFoundError:
        raise Unsupported("legacy format needs LibreOffice, which is not installed")
    except subprocess.TimeoutExpired:
        raise Unsupported(f"LibreOffice conversion timed out after {SOFFICE_TIMEOUT}s")
    except subprocess.CalledProcessError as exc:
        detail = (exc.stderr or b"").decode(errors="replace").strip()[:200]
        raise Unsupported(f"LibreOffice could not convert this file: {detail}")
    stem = os.path.splitext(os.path.basename(path))[0]
    produced = os.path.join(outdir, f"{stem}.{target}")
    if not os.path.exists(produced):
        raise Unsupported("LibreOffice produced no output")
    return produced


def parse(path: str, _depth: int = 0) -> ParsedDoc:
    """Parse any supported document. Raises Unsupported with a reason."""
    if not os.path.isfile(path):
        raise Unsupported(f"not a file: {path}")
    ext = os.path.splitext(path)[1].lower()
    if ext not in HANDLERS and ext not in SOFFICE_TARGETS:
        sniffed = _sniff(path)
        if sniffed is None:
            raise Unsupported(f"unsupported file type '{ext or 'none'}'")
        ext = sniffed

    if ext in HANDLERS:
        doc = HANDLERS[ext](path)
    else:
        if _depth > 1:
            raise Unsupported("nested legacy conversion refused")
        target = SOFFICE_TARGETS[ext]
        converted = _soffice(path, target)
        doc = HANDLERS["." + target](converted)
        doc.converted_via = "libreoffice"
        doc.notes.append(f"converted {ext} -> .{target} via LibreOffice")
        doc.fmt = ext.lstrip(".")

    doc.text = doc.text.replace("\r\n", "\n").replace("\r", "\n")
    if not doc.pages:
        doc.pages = [Page(number=1, text=doc.text, char_start=0)]

    # Tier 3 -- refuse rather than ingest a document with no readable text.
    if len(doc.text.strip()) < MIN_CHARS:
        hint = ("no text layer (scanned image?) -- OCR is out of scope"
                if doc.fmt == "pdf" else "document contains no readable text")
        raise Unsupported(
            f"{hint}: extracted only {len(doc.text.strip())} chars, "
            f"minimum is {MIN_CHARS}")
    return doc


def supported_extensions() -> list[str]:
    return sorted(set(HANDLERS) | set(SOFFICE_TARGETS))


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("supported:", " ".join(supported_extensions()))
        raise SystemExit(0)
    for target in sys.argv[1:]:
        try:
            d = parse(target)
        except Unsupported as exc:
            print(f"REFUSED {target}: {exc}")
            continue
        print(f"OK {target}: fmt={d.fmt} chars={len(d.text)} pages={len(d.pages)}"
              + (f" via={d.converted_via}" if d.converted_via else ""))
        for note in d.notes:
            print(f"   note: {note}")
